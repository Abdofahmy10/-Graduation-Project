import re
from pptx import Presentation
import PyPDF2



def extract_text_from_pptx(file_path):
    try:
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""


def extract_text_from_pdf(file_path):
    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = []
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text.append(page.extract_text())
            return "\n".join(text)
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""



def limit_text_length(text, max_length=512):
    return text[:max_length]



def clean_text(extracted_text):
    cleaned_text = re.sub(r'\s+', ' ', extracted_text)
    return cleaned_text

def chunk_text(text, window_size):
    words = text.split()
    total_words = len(words)
    chunks = []

    for i in range(0, total_words, window_size):
        chunk = ' '.join(words[i:i+window_size])
        chunks.append(chunk)

    return chunks